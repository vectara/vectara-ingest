import json
import os
from core.crawler import Crawler
from omegaconf import OmegaConf
import logging
import io
from datetime import datetime, timedelta
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload
from docx import Document
import pandas as pd
import pptx
from typing import List
import requests

SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]
SERVICE_ACCOUNT_FILE = 'credentials.json'
FILE_UPLOAD_API_URL_TEMPLATE = 'https://api.vectara.io/v2/corpora/{corpus_key}/upload_file'

def get_credentials(delegated_user):
    credentials = service_account.Credentials.from_service_account_file(
        SERVICE_ACCOUNT_FILE, scopes=SCOPES)
    delegated_credentials = credentials.with_subject(delegated_user)
    return delegated_credentials

def download_file(service, file_id):
    try:
        request = service.files().get_media(fileId=file_id)
        file = io.BytesIO()
        downloader = MediaIoBaseDownload(file, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
            print(f"Download {int(status.progress() * 100)}.")
        file.seek(0)  # Reset the file pointer to the beginning
        return file
    except HttpError as error:
        print(f"An error occurred: {error}")
        return None

def export_file(service, file_id, mime_type):
    try:
        request = service.files().export_media(fileId=file_id, mimeType=mime_type)
        file = io.BytesIO()
        downloader = MediaIoBaseDownload(file, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
            print(f"Download {int(status.progress() * 100)}.")
        file.seek(0)  # Reset the file pointer to the beginning
        return file
    except HttpError as error:
        print(f"An error occurred: {error}")
        return None

def extract_text_from_docx(docx_file):
    doc = Document(docx_file)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)

def extract_text_from_csv(csv_file):
    try:
        df = pd.read_csv(csv_file)
        return df.to_string()
    except pd.errors.ParserError as e:
        print(f"Error parsing CSV: {e}")
        return ""

def extract_text_from_xlsx(xlsx_file):
    try:
        dfs = pd.read_excel(xlsx_file, sheet_name=None)
        full_text = []
        for sheet_name, df in dfs.items():
            full_text.append(f"Sheet: {sheet_name}")
            full_text.append(df.to_string())
        return '\n'.join(full_text)
    except Exception as e:
        print(f"Error reading XLSX file: {e}")
        return ""

def extract_text_from_txt(txt_file):
    try:
        return txt_file.read().decode('utf-8')
    except Exception as e:
        print(f"Error reading TXT file: {e}")
        return ""

def extract_text_from_pptx(pptx_file):
    prs = pptx.Presentation(pptx_file)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return '\n'.join(full_text)

def save_local_file(service, file_id, name):
    file_path = os.path.join("/tmp", name)
    try:
        file = download_file(service, file_id)
        if file:
            with open(file_path, 'wb') as f:
                f.write(file.read())
            return file_path
    except Exception as e:
        print(f"Error saving local file: {e}")
    return None

def upload_file_to_vectara(file_path, api_key, corpus_id):
    files = {'file': open(file_path, 'rb')}
    headers = {
        'Authorization': f'Bearer {api_key}',
        'Accept': 'application/json'
    }
    try:
        response = requests.post(FILE_UPLOAD_API_URL_TEMPLATE.format(corpus_key=corpus_id), files=files, headers=headers)
        if response.status_code == 201:
            print("File uploaded successfully")
        else:
            print(f"Failed to upload file: {response.status_code} - {response.text}")
    except Exception as e:
        print(f"Error uploading file: {e}")

class GdriveCrawler(Crawler):

    def __init__(self, cfg: OmegaConf, endpoint: str, customer_id: str, corpus_id: int, api_key: str, delegated_users: List[str]) -> None:
        super().__init__(cfg, endpoint, customer_id, corpus_id, api_key)
        logging.info("Google Drive Crawler initialized")

        self.delegated_users = delegated_users
        self.creds = None
        self.service = None
        self.api_key = api_key
        self.customer_id = customer_id
        self.corpus_id = corpus_id  

    def list_files(self, service, parent_id=None, date_threshold=None):
        results = []
        page_token = None
        query = f"'{parent_id}' in parents and trashed=false and modifiedTime > '{date_threshold}'" if parent_id else f"'root' in parents and trashed=false and modifiedTime > '{date_threshold}'"

        while True:
            try:
                params = {
                    'fields': 'nextPageToken, files(id, name, mimeType, permissions, modifiedTime)',
                    'q': query,
                    'corpora': 'allDrives',
                    'includeItemsFromAllDrives': True,
                    'supportsAllDrives': True
                }
                if page_token:
                    params['pageToken'] = page_token
                response = service.files().list(**params).execute()
                files = response.get('files', [])
                for file in files:
                    permissions = file.get('permissions', [])
                    if any(p.get('displayName') == 'Vectara' or p.get('displayName') == 'all' for p in permissions):
                        results.append(file)
                page_token = response.get('nextPageToken', None)
                if not page_token:
                    break
            except HttpError as error:
                print(f"An error occurred: {error}")
                break
        return results

    def handle_file(self, file):
        file_id = file['id']
        mime_type = file['mimeType']
        name = file['name']
        permissions = file.get('permissions', [])
        
        print(f"Handling file: {name} with MIME type: {mime_type}")

        if not any(p.get('displayName') == 'Vectara' or p.get('displayName') == 'all' for p in permissions):
            print(f"Skipping restricted file: {name}")
            return None

        if mime_type == 'application/vnd.google-apps.document':
            file = export_file(self.service, file_id, 'application/vnd.openxmlformats-officedocument.wordprocessingml.document')
            if file:
                text = extract_text_from_docx(file)
                return text
        elif mime_type == 'application/vnd.google-apps.spreadsheet':
            file = export_file(self.service, file_id, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
            if file:
                text = extract_text_from_xlsx(file)
                return text
        elif mime_type == 'application/vnd.google-apps.presentation':
            file = export_file(self.service, file_id, 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
            if file:
                text = extract_text_from_pptx(file)
                return text
        elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            file = download_file(self.service, file_id)
            if file:
                text = extract_text_from_docx(file)
                return text
        elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            file = download_file(self.service, file_id)
            if file:
                text = extract_text_from_xlsx(file)
                return text
        elif mime_type == 'text/csv':
            file = download_file(self.service, file_id)
            if file:
                text = extract_text_from_csv(file)
                return text
        elif mime_type == 'application/zip':
            return "ZIP files cannot be processed directly."
        elif mime_type == 'application/pdf':
            local_file_path = save_local_file(self.service, file_id, name)
            if local_file_path:
                upload_file_to_vectara(local_file_path, self.api_key, self.corpus_id)
                return f"PDF file '{name}' uploaded."
        elif mime_type == 'application/vnd.google-apps.shortcut':
            return "Shortcuts are not supported."
        elif mime_type == 'application/vnd.google-apps.folder':
            files = self.list_files(self.service, file_id)
            for f in files:
                self.handle_file(f)
            return f"Folder '{name}' with ID '{file_id}'"
        elif mime_type == 'application/vnd.google-apps.form':
            return f"Google Form '{name}' with ID '{file_id}' (Forms content requires Google Forms API)"
        else:
            print(f"Unsupported file type: {mime_type}")
            return None

    def crawl_file(self, file):
        text = self.handle_file(file)
        if text:
            file_id = file['id']
            mime_type = file['mimeType']
            name = file['name']
            created_time = file.get('createdTime', 'N/A')
            modified_time = file.get('modifiedTime', 'N/A')
            owners = ', '.join([owner['displayName'] for owner in file.get('owners', [])])
            size = file.get('size', 'N/A')

            print(f'\n\nCrawling file {name}')

            file_metadata = {
                'id': file_id,
                'mimeType': mime_type,
                'name': name,
                'created_at': created_time,
                'modified_at': modified_time,
                'owners': owners,
                'size': size,
                'source': 'gdrive'
            }

            file_doc = {
                'documentId': f'google-drive-{file_id}',
                'title': name,
                'metadataJson': json.dumps(file_metadata),
                'section': [{
                    'title': name,
                    'text': text,
                }]
            }

            try:
                self.indexer.index_document(file_doc)
            except Exception as e:
                logging.info(f"Error {e} indexing document for file {name}, file_id {file_id}")

    def crawl(self) -> None:
        N = 3  # Number of days to look back
        date_threshold = (datetime.utcnow() - timedelta(days=N)).isoformat() + 'Z'
        
        for user in self.delegated_users:
            print(f"Processing files for user: {user}")
            self.creds = get_credentials(user)
            self.service = build("drive", "v3", credentials=self.creds)
            
            list_files = self.list_files(self.service, date_threshold=date_threshold)
            for file in list_files:
                modified_time = file.get('modifiedTime', 'N/A')
                if modified_time > date_threshold:
                    self.crawl_file(file)
